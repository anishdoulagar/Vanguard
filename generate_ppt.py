"""
CSPM Project Presentation Generator
Generates a comprehensive PowerPoint showcasing the full CSPM platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Palette ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0E, 0x0C, 0x09)   # almost-black
SURFACE  = RGBColor(0x14, 0x12, 0x0F)   # card bg
YELLOW   = RGBColor(0xF5, 0xC5, 0x18)   # electric yellow (primary accent)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)   # cyan
GREEN    = RGBColor(0x4C, 0xAF, 0x7D)   # success
RED      = RGBColor(0xE0, 0x55, 0x55)   # danger
ORANGE   = RGBColor(0xD9, 0x7B, 0x3A)   # warning
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DIM      = RGBColor(0x88, 0x84, 0x7A)   # muted text
BORDER   = RGBColor(0x2A, 0x2D, 0x35)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ─── Helpers ─────────────────────────────────────────────────────────────────

def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, font="Consolas"):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return txb


def label(slide, text, x, y, color=YELLOW, size=7.5):
    """Small all-caps label."""
    return txt(slide, text, x, y, 4, 0.22,
               size=size, bold=True, color=color, font="Calibri")


def bullet_block(slide, items, x, y, w=5.8, size=11, color=DIM, dot_color=YELLOW, line_gap=0.28):
    """Render a list of bullet items."""
    for i, item in enumerate(items):
        # dot
        box(slide, x, y + i * line_gap + 0.08, 0.06, 0.06, fill=dot_color)
        txt(slide, item, x + 0.15, y + i * line_gap, w, 0.28,
            size=size, color=color, font="Calibri")


def pill(slide, text, x, y, fill_color=YELLOW, text_color=BG, w=None, size=8):
    """Colored pill badge."""
    if w is None:
        w = len(text) * 0.08 + 0.25
    b = box(slide, x, y, w, 0.22, fill=fill_color)
    txt(slide, text, x + 0.05, y + 0.02, w, 0.2,
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER, font="Calibri")
    return w


def card(slide, x, y, w, h, title=None, title_color=YELLOW):
    """Draw a card with optional title."""
    box(slide, x, y, w, h, fill=SURFACE, line=BORDER, line_w=Pt(0.75))
    if title:
        txt(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.25,
            size=8.5, bold=True, color=title_color, font="Calibri")
    return (x, y, w, h)


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s


def header_bar(slide, text, sub=None):
    """Standard top section header."""
    box(slide, 0, 0, 13.33, 0.9, fill=SURFACE)
    box(slide, 0, 0.88, 13.33, 0.03, fill=YELLOW)
    txt(slide, text, 0.45, 0.1, 10, 0.55,
        size=22, bold=True, color=WHITE, font="Calibri")
    if sub:
        txt(slide, sub, 0.47, 0.58, 10, 0.3,
            size=10, color=DIM, font="Calibri")


def footer(slide, page_num, total=20):
    box(slide, 0, 7.28, 13.33, 0.22, fill=SURFACE)
    txt(slide, "VANGUARD — CLOUD SECURITY POSTURE MANAGEMENT", 0.4, 7.30, 10, 0.2,
        size=7, color=DIM, font="Calibri")
    txt(slide, f"{page_num} / {total}", 12.5, 7.30, 0.8, 0.2,
        size=7, color=DIM, align=PP_ALIGN.RIGHT, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# background grid lines (decorative)
for i in range(8):
    box(s, 0, i * 0.95, 13.33, 0.005, fill=RGBColor(0x1A, 0x1A, 0x16))
for i in range(15):
    box(s, i * 0.9, 0, 0.005, 7.5, fill=RGBColor(0x1A, 0x1A, 0x16))

# glow slab
box(s, 3.5, 2.2, 6.3, 3.1, fill=RGBColor(0x12, 0x11, 0x0D))

# accent bar
box(s, 3.5, 2.2, 0.07, 3.1, fill=YELLOW)

# title
txt(s, "VANGUARD", 3.75, 2.35, 6, 0.9,
    size=52, bold=True, color=WHITE, font="Calibri")
txt(s, "CLOUD SECURITY POSTURE MANAGEMENT", 3.75, 3.15, 6.5, 0.5,
    size=14, bold=True, color=YELLOW, font="Calibri", align=PP_ALIGN.LEFT)
txt(s, "Multi-cloud security scanning, compliance monitoring, and risk management\n"
       "platform for AWS and Azure cloud environments.",
    3.75, 3.7, 6.2, 0.8, size=11, color=DIM, font="Calibri")

pill(s, "AWS", 3.76, 4.65, fill_color=ORANGE, text_color=WHITE, w=0.55, size=9)
pill(s, "AZURE", 3.76+0.65, 4.65, fill_color=RGBColor(0x00,0x78,0xD4), text_color=WHITE, w=0.65, size=9)
pill(s, "137+ RULES", 3.76+0.65+0.75, 4.65, fill_color=GREEN, text_color=WHITE, w=0.95, size=9)
pill(s, "6 FRAMEWORKS", 3.76+0.65+0.75+1.05, 4.65, fill_color=CYAN, text_color=BG, w=1.1, size=9)

txt(s, "Anish Doulagar  ·  2026", 3.76, 5.15, 6, 0.3, size=9, color=DIM, font="Calibri")

# shield icon (text-art)
txt(s, "🛡", 0.6, 2.8, 2, 1.5, size=72, color=YELLOW, align=PP_ALIGN.CENTER)

footer(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "AGENDA", "What we'll cover")

AGENDA = [
    ("01", "Problem & Motivation",      "Why cloud security posture matters"),
    ("02", "Architecture Overview",     "System design and technology stack"),
    ("03", "Authentication & Security", "JWT, MFA, RBAC, team isolation"),
    ("04", "Cloud Connectors",          "AWS & Azure service collectors"),
    ("05", "Policy Engine",             "137+ rules across 6 frameworks"),
    ("06", "Risk Scoring",              "Penalty model and severity levels"),
    ("07", "Dashboard & UI",            "Live metrics and interactive views"),
    ("08", "Accounts Management",       "Multi-cloud account lifecycle"),
    ("09", "Scan Engine",               "On-demand and scheduled scans"),
    ("10", "Findings & Compliance",     "Results, tracking, and remediation"),
    ("11", "Alerting System",           "Email alerts with HTML reports"),
    ("12", "Custom Policies",           "User-defined rules and operators"),
    ("13", "Audit Log",                 "Full action trail across the platform"),
    ("14", "Team & RBAC System",        "Multi-team isolation and permissions"),
    ("15", "Deployment",                "Docker Compose, Nginx, prod config"),
]

col_w = 5.8
for i, (num, title, sub) in enumerate(AGENDA):
    col = i // 8
    row = i % 8
    cx = 0.55 + col * (col_w + 0.65)
    cy = 1.1 + row * 0.74

    box(s, cx, cy, 0.38, 0.38, fill=YELLOW)
    txt(s, num, cx, cy + 0.04, 0.38, 0.3, size=10, bold=True,
        color=BG, align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, title, cx + 0.5, cy + 0.01, col_w - 0.55, 0.26,
        size=11, bold=True, color=WHITE, font="Calibri")
    txt(s, sub, cx + 0.5, cy + 0.24, col_w - 0.55, 0.2,
        size=8.5, color=DIM, font="Calibri")

footer(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM & MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "PROBLEM & MOTIVATION", "The cloud security challenge")

# Big stat cards
stats = [
    ("$4.9M",   "Average cost of a cloud\nsecurity breach (IBM 2024)", YELLOW),
    ("82%",     "Of breaches involve cloud\nmisconfigurations", RED),
    ("200+ days","Mean time to detect\na cloud breach", ORANGE),
    ("10,000+", "New misconfigurations\ncreated per day in enterprise", CYAN),
]
for i, (stat, label_txt, col) in enumerate(stats):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.1, 2.9, 1.5)
    txt(s, stat, cx + 0.2, 1.2, 2.5, 0.7, size=28, bold=True, color=col, font="Calibri")
    txt(s, label_txt, cx + 0.2, 1.85, 2.5, 0.65, size=9, color=DIM, font="Calibri")

# Problem bullets
card(s, 0.4, 2.85, 5.9, 3.0, title="THE PROBLEM")
probs = [
    "Cloud environments grow faster than security teams can monitor",
    "Misconfigurations are the #1 cause of cloud security incidents",
    "Manual auditing is slow, error-prone, and doesn't scale",
    "Compliance frameworks (CIS, PCI, HIPAA) require continuous validation",
    "Multi-cloud visibility across AWS & Azure is fragmented",
    "No centralized view of organizational security posture",
]
bullet_block(s, probs, 0.6, 3.2, w=5.5, size=10.5)

# Solution bullets
card(s, 6.6, 2.85, 6.3, 3.0, title="THE SOLUTION — VANGUARD")
sols = [
    "Automated scanning of cloud resources against 137+ security rules",
    "Unified dashboard with real-time risk scores and trending",
    "Multi-framework compliance mapping (CIS, NIST, PCI, HIPAA, SOC2, ISO)",
    "Team-based access control — isolated views per organization",
    "Custom policy engine for org-specific security requirements",
    "Email alerting and scheduled background scans",
]
bullet_block(s, sols, 6.8, 3.2, w=5.9, size=10.5, dot_color=GREEN)

footer(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ARCHITECTURE OVERVIEW", "System design and data flow")

# Layers
layers = [
    ("FRONTEND",  "React 18 + Vite",         "Dashboard · Accounts · Policies · Audit · Teams", CYAN,   0.9),
    ("API",       "FastAPI (Python 3.11)",    "REST endpoints · JWT auth · RBAC middleware",      YELLOW, 2.25),
    ("ENGINE",    "Scan Engine + Scheduler",  "Connectors · Policy Evaluator · Risk Scorer · APScheduler", GREEN, 3.6),
    ("DATABASE",  "PostgreSQL 15",            "users · cloud_accounts · scan_results · audit_log · teams", ORANGE, 4.95),
    ("INFRA",     "Docker Compose + Nginx",   "Containerized services · TLS termination · reverse proxy",  DIM,    6.3),
]

for (layer, tech, detail, col, y) in layers:
    box(s, 0.4, y, 0.08, 0.7, fill=col)
    card(s, 0.6, y, 12.3, 0.72)
    txt(s, layer, 0.85, y + 0.06, 1.3, 0.28, size=9, bold=True, color=col, font="Calibri")
    txt(s, tech,  2.25, y + 0.06, 4.5, 0.28, size=11, bold=True, color=WHITE, font="Calibri")
    txt(s, detail, 2.25, y + 0.36, 10.4, 0.28, size=9.5, color=DIM, font="Calibri")

    # arrows between layers (except last)
    if y < 6.0:
        txt(s, "↓", 6.5, y + 0.7, 0.4, 0.3, size=14, color=BORDER, align=PP_ALIGN.CENTER)

footer(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "TECHNOLOGY STACK", "Libraries, frameworks, and services")

stack = [
    ("Backend",      YELLOW, [
        "Python 3.11  /  FastAPI (async/await)",
        "asyncpg — async PostgreSQL driver",
        "Pydantic v2 — request validation",
        "python-jose — JWT HS256 tokens",
        "bcrypt — password hashing",
        "cryptography (Fernet) — credential encryption",
        "APScheduler — background job scheduling",
        "pyotp — TOTP MFA (RFC 6238)",
    ]),
    ("Frontend",     CYAN, [
        "React 18 with Vite build tooling",
        "Recharts — scan trend line charts",
        "Inline CSS with CSS variables (dark/light theme)",
        "Custom icon components (SVG)",
        "No external UI library — fully custom",
    ]),
    ("Cloud SDKs",   GREEN, [
        "boto3 — AWS SDK (23+ service clients)",
        "azure-identity — Azure credential chain",
        "azure-mgmt-* — ARM management clients",
        "azure-mgmt-resource / subscription",
    ]),
    ("Infrastructure", ORANGE, [
        "PostgreSQL 15 (Docker)",
        "Docker Compose (dev + prod configs)",
        "Nginx — reverse proxy + TLS",
        "SMTP email (Gmail / SendGrid)",
        "uvicorn — ASGI server",
    ]),
]

for i, (cat, col, items) in enumerate(stack):
    col_x = 0.4 + i * 3.2
    card(s, col_x, 1.05, 3.0, 5.8, title=cat, title_color=col)
    bullet_block(s, items, col_x + 0.18, 1.45, w=2.65, size=9.5, dot_color=col)

footer(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — AUTHENTICATION & SECURITY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "AUTHENTICATION & SECURITY", "Multi-layer access control")

# Auth flow diagram boxes
flow = [
    (0.5,  2.0, "USER LOGIN\nEmail + Password"),
    (3.3,  2.0, "JWT TOKEN\nHS256, 24h TTL"),
    (6.1,  2.0, "MFA CHECK\nTOTP (RFC 6238)"),
    (8.9,  2.0, "RBAC GUARD\nrole + team check"),
    (11.5, 2.0, "RESOURCE\nACCESS"),
]
for (x, y, label_txt) in flow:
    card(s, x, y, 2.6, 0.85)
    txt(s, label_txt, x+0.15, y+0.1, 2.3, 0.7, size=9, bold=True,
        color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)

# arrows
for xi in [3.05, 5.85, 8.65, 11.25]:
    txt(s, "→", xi, 2.3, 0.3, 0.3, size=16, color=YELLOW, align=PP_ALIGN.CENTER)

# Feature cards
feats = [
    ("JWT Authentication", YELLOW, [
        "HS256-signed tokens, 24-hour expiry",
        "Bearer token on every authenticated request",
        "Token payload: user ID, email, role",
        "/auth/me validates and refreshes session",
    ]),
    ("MFA — TOTP", CYAN, [
        "RFC 6238 time-based one-time passwords",
        "QR code provisioning (pyotp + qrcode)",
        "6-digit codes, 30-second window",
        "Required after password login if enabled",
        "Per-user opt-in, admin can enforce",
    ]),
    ("Password Security", GREEN, [
        "bcrypt with cost factor 12",
        "Forgot-password email with signed tokens",
        "Admin-initiated force-reset",
        "Token single-use, 1-hour TTL",
    ]),
    ("Role-Based Access", ORANGE, [
        "4 roles: superadmin, admin, analyst, viewer",
        "require_role() dependency on every endpoint",
        "Role hierarchy: cannot escalate own role",
        "Admin can only manage analyst/viewer",
        "Fernet-encrypted cloud credentials at rest",
    ]),
]

for i, (title, col, items) in enumerate(feats):
    cx = 0.4 + i * 3.2
    card(s, cx, 3.15, 3.05, 3.7, title=title, title_color=col)
    bullet_block(s, items, cx + 0.18, 3.5, w=2.7, size=9.5, dot_color=col)

footer(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — CLOUD CONNECTORS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "CLOUD CONNECTORS", "AWS & Azure service collectors")

# AWS column
card(s, 0.4, 1.05, 6.0, 5.85, title="AWS CONNECTOR", title_color=ORANGE)
aws_services = [
    "IAM — Users, roles, policies, access keys, password policy",
    "S3 — Bucket ACLs, encryption, versioning, public access blocks",
    "EC2 — Instances, security groups, EBS volumes, VPC config",
    "RDS — DB instances, snapshots, encryption, public access",
    "Lambda — Functions, environment variables, runtime versions",
    "CloudTrail — Trail configs, log validation, S3 destinations",
    "KMS — Key metadata, rotation status, key policies",
    "EKS — Cluster configs, endpoint access, logging",
    "ECR — Repositories, scan-on-push, image tag mutability",
    "CloudWatch — Log groups, metric alarms, retention policies",
    "VPC — Flow logs, default VPC usage, NACL/SG rules",
    "ElastiCache — Cluster encryption, auth token configs",
    "SNS — Topic encryption, cross-account access policies",
    "SQS — Queue encryption, visibility timeout, dead-letter",
]
bullet_block(s, aws_services, 0.6, 1.45, w=5.6, size=9.5, dot_color=ORANGE)

# Azure column
card(s, 6.7, 1.05, 6.25, 5.85, title="AZURE CONNECTOR", title_color=CYAN)
azure_services = [
    "Storage Accounts — Encryption, HTTPS, public blob access",
    "Virtual Machines — Managed disks, JIT access, extensions",
    "Network Security Groups — Inbound rules, open ports",
    "Key Vaults — Soft delete, purge protection, RBAC",
    "SQL Databases — TDE, auditing, threat detection",
    "AKS — Kubernetes version, RBAC, network policies",
    "App Service — HTTPS only, TLS version, auth settings",
    "MySQL / PostgreSQL — SSL enforcement, backup retention",
    "Event Hub — Namespace encryption, network rules",
    "Service Bus — CMK encryption, network isolation",
    "Azure Monitor — Diagnostic settings, log retention",
    "Logic Apps — Auth configurations, connection security",
    "Azure Functions — Auth level, managed identity",
    "Defender for Cloud — Plans per resource type",
    "Network Watchers — Coverage per region",
]
bullet_block(s, azure_services, 6.9, 1.45, w=5.85, size=9.5, dot_color=CYAN)

footer(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — POLICY ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "POLICY ENGINE", "137+ rules across 6 compliance frameworks")

# Framework pills row
frameworks = [
    ("CIS Benchmarks", YELLOW),
    ("NIST 800-53", CYAN),
    ("PCI-DSS v4", RED),
    ("HIPAA", GREEN),
    ("SOC 2", ORANGE),
    ("ISO 27001", DIM),
    ("GDPR", RGBColor(0x60, 0x8B, 0xCC)),
]
px = 0.45
for name, col in frameworks:
    w = pill(s, name, px, 1.08, fill_color=col,
             text_color=BG if col != DIM else WHITE, size=8)
    px += w + 0.2

# Rule anatomy card
card(s, 0.4, 1.55, 5.5, 2.8, title="RULE STRUCTURE")
txt(s, '{\n  "id":          "AWS-IAM-001",\n  "title":       "Root account access key",\n  "severity":    "CRITICAL",\n  "frameworks":  ["CIS", "PCI", "NIST"],\n  "resource_type": "iam_user",\n  "description": "Root should not have active keys",\n  "remediation": "Delete root account access keys"\n}',
    0.58, 1.85, 5.1, 2.35, size=8.5, color=GREEN, font="Consolas")

# Rule counts by category
card(s, 6.15, 1.55, 6.75, 2.8, title="RULE COVERAGE BY SERVICE")
rules = [
    ("IAM",          "22 rules", YELLOW),
    ("S3",           "14 rules", CYAN),
    ("EC2 / VPC",    "18 rules", GREEN),
    ("RDS / Aurora", "10 rules", ORANGE),
    ("Lambda",       " 6 rules", DIM),
    ("CloudTrail",   " 8 rules", YELLOW),
    ("KMS",          " 5 rules", CYAN),
    ("Azure Storage","12 rules", RGBColor(0x00,0x78,0xD4)),
    ("Azure NSG",    " 9 rules", GREEN),
    ("Azure AD/KV",  "16 rules", ORANGE),
    ("+ 30 more…",   "137 total", DIM),
]
for i, (svc, count, col) in enumerate(rules):
    row_y = 1.9 + i * 0.29
    txt(s, svc,   6.3,  row_y, 3.0, 0.28, size=9.5, color=DIM, font="Calibri")
    txt(s, count, 9.8,  row_y, 1.5, 0.28, size=9.5, bold=True, color=col, font="Calibri")

# Severity breakdown
card(s, 0.4, 4.55, 12.5, 2.35, title="SEVERITY LEVELS & PENALTY MODEL")
sevs = [
    ("CRITICAL", "10 pts", RED,    "Root exposure, public S3 data, unencrypted secrets, open RDP/SSH globally"),
    ("HIGH",     " 7 pts", ORANGE, "MFA disabled, CloudTrail off, unpatched VMs, public DB endpoints"),
    ("MEDIUM",   " 4 pts", YELLOW, "Missing tags, weak TLS, no backup retention, logging gaps"),
    ("LOW",      " 1 pt",  GREEN,  "Minor config issues, best-practice deviations, informational"),
]
for i, (sev, pts, col, example) in enumerate(sevs):
    sx = 0.6 + i * 3.1
    pill(s, sev, sx, 4.88, fill_color=col, text_color=BG if col!=GREEN else BG, w=1.0, size=8)
    txt(s, pts,     sx + 1.1, 4.88, 0.7, 0.22, size=9, bold=True, color=col, font="Calibri")
    txt(s, example, sx,       5.2,  2.9, 0.55, size=8.5, color=DIM, font="Calibri")

footer(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — DASHBOARD & UI
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "DASHBOARD & UI", "Live metrics, charts, and interactive views")

# Pages list
card(s, 0.4, 1.05, 3.3, 5.85, title="PAGES / VIEWS")
pages = [
    ("Dashboard",      "Overview metrics + trend chart"),
    ("Accounts",       "Cloud account management"),
    ("Quick Scan",     "Ad-hoc scan with live results"),
    ("History",        "Past scan results & findings"),
    ("Alerts",         "Alert config + notification log"),
    ("Policies",       "Built-in & custom rule browser"),
    ("Audit Log",      "Full action trail"),
    ("Teams",          "Team CRUD + member management"),
    ("My Team",        "Admin's team member view"),
    ("Users",          "Platform user management"),
]
for i, (pg, desc) in enumerate(pages):
    y = 1.45 + i * 0.54
    txt(s, pg,   0.6, y,        2.1, 0.27, size=10, bold=True, color=YELLOW, font="Calibri")
    txt(s, desc, 0.6, y + 0.25, 3.0, 0.22, size=8.5, color=DIM, font="Calibri")

# Dashboard features
card(s, 3.95, 1.05, 4.8, 2.7, title="DASHBOARD METRICS")
metrics = [
    "Total cloud accounts connected",
    "Accounts scanned vs. unscanned",
    "Last 90-day trend chart (Recharts)",
    "Findings by severity: CRITICAL / HIGH / MEDIUM / LOW",
    "Overall risk score (0–100, penalty model)",
    "Recently flagged findings with quick links",
]
bullet_block(s, metrics, 4.1, 1.45, w=4.4, size=9.5, dot_color=CYAN)

card(s, 3.95, 3.95, 4.8, 2.95, title="THEMING")
themes = [
    "Dark mode: cyberpunk electric-yellow palette",
    "Light mode: high-contrast warm-beige theme",
    "CSS custom properties (--bg, --surface, --accent…)",
    "Single toggle — all components respond instantly",
    "Monospace (Consolas) + UI (Calibri) fonts",
    "Vanguard brand: shield icon + electric accents",
]
bullet_block(s, themes, 4.1, 4.35, w=4.4, size=9.5, dot_color=YELLOW)

# UX features
card(s, 9.0, 1.05, 3.9, 5.85, title="UX FEATURES")
ux = [
    "Top horizontal navigation bar",
    "Role-filtered nav items (minRole)",
    "ADMIN ▾ dropdown (Teams / Users)",
    "Live scan progress indicator",
    "Findings expandable per-resource",
    "Acknowledge / Resolve finding status",
    "Search + filter on all data tables",
    "Compliance framework filter tabs",
    "Dark/light theme toggle",
    "MFA security indicator in nav",
    "Per-user avatar + role badge",
    "Empty states with onboarding hints",
    "Responsive overflow-x tables",
]
bullet_block(s, ux, 9.18, 1.45, w=3.5, size=9.5, dot_color=GREEN)

footer(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — ACCOUNTS MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "CLOUD ACCOUNTS MANAGEMENT", "Lifecycle of connected cloud accounts")

# Lifecycle flow
steps = [
    ("ADD ACCOUNT",   "Enter cloud credentials\n(AWS keys or Azure SP)"),
    ("ENCRYPT",       "Fernet-encrypt creds\nstored in PostgreSQL"),
    ("TEST CONN",     "Validate credentials\nbefore saving"),
    ("ASSIGN TEAM",   "Auto/manual team\nassignment on create"),
    ("SCAN",          "Run policy engine\nagainst live resources"),
    ("VIEW RESULTS",  "Findings, risk score,\ncompliance report"),
]
for i, (step, desc) in enumerate(steps):
    cx = 0.4 + i * 2.15
    col = [YELLOW, GREEN, CYAN, ORANGE, GREEN, YELLOW][i]
    card(s, cx, 1.1, 2.0, 1.3)
    txt(s, step, cx+0.12, 1.2, 1.78, 0.3, size=9, bold=True, color=col, font="Calibri")
    txt(s, desc, cx+0.12, 1.55, 1.78, 0.7, size=8.5, color=DIM, font="Calibri")
    if i < 5:
        txt(s, "→", cx+2.02, 1.6, 0.22, 0.3, size=13, color=BORDER, align=PP_ALIGN.CENTER)

# AWS creds card
card(s, 0.4, 2.65, 4.0, 2.55, title="AWS CREDENTIALS")
aws_creds = [
    "Access Key ID + Secret Access Key",
    "Optional: Session Token (temporary creds)",
    "Optional: explicit AWS Region",
    "Validated via IAM GetCallerIdentity",
    "Stored encrypted — never returned to client",
]
bullet_block(s, aws_creds, 0.6, 3.05, w=3.6, size=10, dot_color=ORANGE)

# Azure creds card
card(s, 4.65, 2.65, 4.0, 2.55, title="AZURE CREDENTIALS")
az_creds = [
    "Tenant ID + Client ID + Client Secret",
    "Subscription ID for resource scoping",
    "Service Principal with Reader role minimum",
    "Validated via subscription list call",
    "ChainedTokenCredential support",
]
bullet_block(s, az_creds, 4.85, 3.05, w=3.6, size=10, dot_color=CYAN)

# Operations card
card(s, 8.9, 2.65, 4.0, 2.55, title="ACCOUNT OPERATIONS")
ops = [
    "Edit credentials (re-encrypt + re-validate)",
    "Delete account (cascades scan history)",
    "Test connection without scanning",
    "Move to different team (superadmin/admin)",
    "Scan-all: bulk scan all team accounts",
    "Schedule: auto-scan every N hours",
]
bullet_block(s, ops, 9.1, 3.05, w=3.6, size=10, dot_color=GREEN)

# Security note
box(s, 0.4, 5.4, 12.5, 0.7, fill=RGBColor(0x1A, 0x14, 0x08))
box(s, 0.4, 5.4, 0.06, 0.7, fill=YELLOW)
txt(s, "SECURITY:  Cloud credentials are encrypted with Fernet (AES-128-CBC + HMAC-SHA256) before "
       "database storage. The encryption key is stored in the server environment, never in the database. "
       "Credentials are decrypted only at scan runtime and never exposed through any API response.",
    0.6, 5.5, 12.1, 0.5, size=9, color=DIM, font="Calibri")

footer(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — SCAN ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "SCAN ENGINE", "On-demand and scheduled cloud scanning")

# Scan pipeline
pipeline = [
    ("1. LOAD ACCOUNT",  "Decrypt creds from DB\nDetermine cloud provider", YELLOW),
    ("2. COLLECT",       "Call cloud SDK APIs\nNormalize to CloudResource", CYAN),
    ("3. EVALUATE",      "Run all matching rules\nper resource_type", GREEN),
    ("4. SCORE",         "Penalty model\nCRIT/HIGH/MED/LOW → 0–100", ORANGE),
    ("5. PERSIST",       "Store scan_result row\nwith JSON findings", RED),
    ("6. NOTIFY",        "Trigger email alert\nif threshold exceeded", DIM),
]
for i, (step, desc, col) in enumerate(pipeline):
    cx = 0.4 + i * 2.15
    box(s, cx, 1.1, 2.0, 0.3, fill=col)
    txt(s, step, cx+0.08, 1.13, 1.85, 0.25, size=8, bold=True,
        color=BG, align=PP_ALIGN.CENTER, font="Calibri")
    card(s, cx, 1.38, 2.0, 1.0)
    txt(s, desc, cx+0.12, 1.45, 1.78, 0.85, size=9, color=DIM, font="Calibri")
    if i < 5:
        txt(s, "→", cx+2.0, 1.65, 0.22, 0.3, size=13, color=BORDER, align=PP_ALIGN.CENTER)

# Scan types
card(s, 0.4, 2.65, 5.9, 4.25, title="SCAN TYPES")
scan_types = [
    ("On-demand scan",       "User triggers scan on a specific saved account.\n"
                             "Results stream in real-time. Shown on results page."),
    ("Ad-hoc quick scan",    "Enter creds inline (not saved) for instant scan.\n"
                             "Great for one-off audits without saving credentials."),
    ("Scheduled scan",       "APScheduler job runs every N hours per account.\n"
                             "Configurable interval. Results stored automatically."),
    ("Bulk scan-all",        "Admin/analyst triggers scan on all team accounts\n"
                             "in parallel. Status tracked with scanning_accounts set."),
]
for i, (stype, sdesc) in enumerate(scan_types):
    sy = 3.05 + i * 0.95
    txt(s, stype, 0.6, sy, 2.5, 0.28, size=10, bold=True, color=YELLOW, font="Calibri")
    txt(s, sdesc, 0.6, sy+0.28, 5.5, 0.55, size=9, color=DIM, font="Calibri")

# Concurrency / performance
card(s, 6.55, 2.65, 6.35, 4.25, title="CONCURRENCY & PERFORMANCE")
perf = [
    "asyncio + asyncpg — non-blocking DB calls throughout",
    "Scanning runs in ThreadPoolExecutor (boto3 is sync)",
    "Per-account scan lock prevents duplicate concurrent scans",
    "scanning_accounts in-memory set for live status polling",
    "Findings stored as compressed JSON in scan_results",
    "Resource normalization decouples connectors from rules",
    "Rule evaluation: O(rules × resources) per scan",
    "All cloud SDK calls wrapped in try/except per service",
    "Single service failure never aborts the whole scan",
    "Scan metadata: duration, resource count, finding counts",
]
bullet_block(s, perf, 6.73, 3.05, w=5.95, size=9.5, dot_color=CYAN)

footer(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — FINDINGS & COMPLIANCE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "FINDINGS & COMPLIANCE", "Results, tracking, and remediation guidance")

# Finding anatomy
card(s, 0.4, 1.05, 5.7, 3.5, title="FINDING STRUCTURE")
txt(s,
    '{\n'
    '  "rule_id":       "AWS-S3-003",\n'
    '  "resource_id":   "my-prod-bucket",\n'
    '  "resource_type": "s3_bucket",\n'
    '  "severity":      "CRITICAL",\n'
    '  "title":         "S3 bucket is publicly accessible",\n'
    '  "description":   "Bucket has public ACL or Block disabled",\n'
    '  "remediation":   "Enable Block Public Access on bucket",\n'
    '  "frameworks":    ["CIS-2.1.5", "PCI-3.4", "HIPAA"],\n'
    '  "status":        "acknowledged",\n'
    '  "details":       { "acl": "public-read", "region": "us-east-1" }\n'
    '}',
    0.58, 1.42, 5.3, 2.9, size=8, color=GREEN, font="Consolas")

# Finding statuses
card(s, 6.35, 1.05, 6.55, 1.65, title="FINDING STATUS TRACKING")
statuses = [
    ("OPEN",         "New finding, not yet addressed", RED),
    ("ACKNOWLEDGED", "Noted — tracking for remediation", ORANGE),
    ("RESOLVED",     "Fixed and confirmed by re-scan", GREEN),
]
for i, (st, desc, col) in enumerate(statuses):
    sy = 1.42 + i * 0.4
    pill(s, st, 6.5, sy, fill_color=col, text_color=WHITE if col==RED else BG, w=1.2, size=7.5)
    txt(s, desc, 7.85, sy, 4.8, 0.28, size=9, color=DIM, font="Calibri")

# Compliance frameworks
card(s, 6.35, 2.9, 6.55, 3.65, title="COMPLIANCE FRAMEWORK MAPPING")
fw_detail = [
    ("CIS Benchmarks",  "Industry standard hardening guidelines for AWS & Azure"),
    ("NIST 800-53",     "Federal security controls — mapped per control family"),
    ("PCI-DSS v4",      "Payment card data security — 12 requirements mapped"),
    ("HIPAA",           "Healthcare data protection — administrative & technical"),
    ("SOC 2",           "Service org security — Trust Services Criteria"),
    ("ISO 27001",       "Information security management — Annex A controls"),
    ("GDPR",            "Data protection — privacy by design principles"),
]
for i, (fw, desc) in enumerate(fw_detail):
    fy = 3.25 + i * 0.44
    txt(s, fw,   6.5,  fy, 2.0, 0.28, size=9, bold=True, color=YELLOW, font="Calibri")
    txt(s, desc, 8.65, fy, 3.9, 0.28, size=9,            color=DIM,    font="Calibri")

# Risk score explanation
card(s, 0.4, 4.75, 12.5, 1.55, title="RISK SCORE MODEL")
txt(s, "Score = 100 − Σ(severity_penalty × finding_count)   ·   clamped to [0, 100]",
    0.6, 5.1, 8.0, 0.32, size=11, bold=True, color=YELLOW, font="Consolas")
score_bands = [
    ("90–100", "EXCELLENT", GREEN),
    ("70–89",  "GOOD",      CYAN),
    ("50–69",  "FAIR",      YELLOW),
    ("30–49",  "POOR",      ORANGE),
    ("0–29",   "CRITICAL",  RED),
]
for i, (band, label_t, col) in enumerate(score_bands):
    bx = 0.6 + i * 2.4
    txt(s, band,    bx, 5.5, 1.5, 0.25, size=10, bold=True, color=col, font="Calibri")
    txt(s, label_t, bx, 5.75, 1.5, 0.25, size=8.5, color=DIM, font="Calibri")

footer(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — ALERTING SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ALERTING SYSTEM", "Email notifications with HTML security reports")

card(s, 0.4, 1.05, 6.0, 5.85, title="ALERT CONFIGURATION")
alert_items = [
    ("Trigger Modes", [
        "Score drops below custom threshold (e.g., < 70)",
        "New CRITICAL finding discovered in scan",
        "Both score threshold AND critical finding",
    ]),
    ("Scope Options", [
        "Per-account alert settings (individual thresholds)",
        "System-wide default alert settings (superadmin)",
        "Users set per-account overrides independently",
    ]),
    ("Email Content", [
        "HTML report with Vanguard branding",
        "Account name, scan timestamp, risk score",
        "Full finding list grouped by severity",
        "Remediation steps per finding",
        "Compliance framework tags per finding",
    ]),
    ("Delivery", [
        "SMTP via Gmail or SendGrid",
        "Configurable sender/recipients per account",
        "Email sent in background thread (non-blocking)",
        "Alert history stored in alert_history table",
        "Deduplication: no re-alert if score unchanged",
    ]),
]
y_off = 1.45
for section, items in alert_items:
    txt(s, section, 0.6, y_off, 5.5, 0.28, size=10, bold=True, color=YELLOW, font="Calibri")
    bullet_block(s, items, 0.6, y_off + 0.28, w=5.5, size=9.5)
    y_off += 0.28 + len(items) * 0.28 + 0.18

# Alert history card
card(s, 6.65, 1.05, 6.25, 2.85, title="ALERT HISTORY LOG")
hist = [
    "Every alert email stored in alert_history table",
    "Fields: account, trigger type, score at time, timestamp",
    "Displayed in Alerts page as a scrollable feed",
    "Admin can review alert frequency per account",
    "Useful for compliance evidence and audit trails",
]
bullet_block(s, hist, 6.85, 1.45, w=5.8, size=9.5, dot_color=ORANGE)

# Scheduler card
card(s, 6.65, 4.1, 6.25, 2.8, title="SCHEDULED SCANNING")
sched = [
    "APScheduler BackgroundScheduler (async-safe)",
    "Per-account cron interval (hours configurable)",
    "Jobs stored in scheduled_jobs DB table",
    "On startup: all active jobs are re-registered",
    "Scan completion updates next_run_at in DB",
    "Admin can enable/disable/update schedule per account",
]
bullet_block(s, sched, 6.85, 4.5, w=5.8, size=9.5, dot_color=CYAN)

footer(s, 13)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — CUSTOM POLICIES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "CUSTOM POLICIES", "User-defined security rules with JSON operators")

card(s, 0.4, 1.05, 5.7, 5.85, title="CUSTOM RULE ANATOMY")
txt(s,
    '{\n'
    '  "id":           "CUSTOM-001",\n'
    '  "title":        "Prod buckets must have versioning",\n'
    '  "description":  "All prod- prefixed S3 buckets...",\n'
    '  "severity":     "HIGH",\n'
    '  "resource_type":"s3_bucket",\n'
    '  "condition": {\n'
    '    "field":    "tags.Environment",\n'
    '    "operator": "equals",\n'
    '    "value":    "production"\n'
    '  },\n'
    '  "frameworks":   ["internal-policy"],\n'
    '  "remediation":  "Enable S3 versioning on bucket"\n'
    '}',
    0.58, 1.42, 5.3, 3.85, size=9, color=GREEN, font="Consolas")

txt(s, "Stored as JSON in custom_rules.json · Evaluated alongside built-in rules",
    0.58, 5.38, 5.3, 0.3, size=8.5, color=DIM, font="Calibri")

# Operators
card(s, 6.35, 1.05, 6.55, 2.8, title="SUPPORTED OPERATORS")
ops_list = [
    ("equals",           "Exact match on field value"),
    ("not_equals",       "Field value is not X"),
    ("contains",         "String contains substring"),
    ("not_contains",     "String does not contain"),
    ("exists",           "Field is present and not null"),
    ("not_exists",       "Field is absent or null"),
    ("greater_than",     "Numeric comparison >"),
    ("less_than",        "Numeric comparison <"),
    ("in",               "Value is in a list"),
    ("not_in",           "Value is not in a list"),
]
for i, (op, desc) in enumerate(ops_list):
    oy = 1.42 + i * 0.33
    txt(s, op,   6.5,  oy, 2.0, 0.28, size=9,   bold=True, color=CYAN, font="Consolas")
    txt(s, desc, 8.7,  oy, 3.9, 0.28, size=9,              color=DIM,  font="Calibri")

# CRUD operations
card(s, 6.35, 4.05, 6.55, 2.85, title="CUSTOM RULE MANAGEMENT")
crud = [
    "Create rule via Policies page form (analyst+)",
    "Full validation on condition structure",
    "Edit any existing custom rule (own or team)",
    "Delete custom rule with confirmation",
    "Rules evaluated on every scan automatically",
    "Rules appear in built-in policy browser alongside defaults",
    "Findings tagged with framework 'custom'",
    "Risk score includes custom rule violations",
]
bullet_block(s, crud, 6.55, 4.45, w=5.9, size=9.5, dot_color=YELLOW)

footer(s, 14)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — AUDIT LOG
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "AUDIT LOG", "Immutable record of all significant platform actions")

card(s, 0.4, 1.05, 5.7, 5.85, title="LOGGED ACTIONS")
logged = [
    ("login",               "User signed in", GREEN),
    ("signup",              "New user registered", CYAN),
    ("create_account",      "Cloud account added", CYAN),
    ("delete_account",      "Cloud account removed", RED),
    ("scan_account",        "Scan triggered on account", GREEN),
    ("bulk_scan",           "Scan-all triggered", GREEN),
    ("delete_user",         "User deleted by admin", RED),
    ("change_user_role",    "Role changed by admin", ORANGE),
    ("create_custom_rule",  "New custom rule added", CYAN),
    ("delete_custom_rule",  "Custom rule removed", RED),
    ("update_alert_settings","Alert config changed", YELLOW),
    ("create_team",         "New team created", CYAN),
    ("add_team_member",     "User added to team", GREEN),
    ("remove_team_member",  "User removed from team", RED),
    ("mfa_enabled",         "MFA activated by user", GREEN),
    ("password_reset",      "Password reset by user/admin", ORANGE),
]
for i, (action, desc, col) in enumerate(logged):
    ay = 1.42 + i * 0.33
    pill(s, action, 0.55, ay, fill_color=col,
         text_color=BG if col != RED else WHITE, w=1.85, size=6.5)
    txt(s, desc, 2.55, ay, 3.4, 0.28, size=9, color=DIM, font="Calibri")

card(s, 6.35, 1.05, 6.55, 2.6, title="LOG ENTRY FIELDS")
fields = [
    ("user_id",       "UUID of the acting user"),
    ("user_email",    "Email for display"),
    ("action",        "Action type key"),
    ("resource_id",   "Target resource UUID"),
    ("resource_name", "Human-readable resource"),
    ("ip_address",    "Requester IP (from headers)"),
    ("detail",        "JSON blob with context"),
    ("created_at",    "UTC timestamp"),
]
for i, (field, desc) in enumerate(fields):
    fy = 1.42 + i * 0.28
    txt(s, field, 6.5,  fy, 2.0, 0.26, size=8.5, bold=True, color=CYAN, font="Consolas")
    txt(s, desc,  8.65, fy, 4.0, 0.26, size=8.5,            color=DIM,  font="Calibri")

card(s, 6.35, 3.85, 6.55, 3.05, title="VISIBILITY BY ROLE")
vis = [
    ("superadmin", "ALL logs — complete platform audit trail",   YELLOW),
    ("admin",      "All actions by users in their teams",        CYAN),
    ("analyst",    "Their own actions only",                     GREEN),
    ("viewer",     "No access to audit log",                     DIM),
]
for i, (role, desc, col) in enumerate(vis):
    ry = 4.2 + i * 0.6
    pill(s, role.upper(), 6.5, ry, fill_color=col, text_color=BG if col!=DIM else WHITE, w=1.2, size=7.5)
    txt(s, desc, 7.85, ry, 4.9, 0.42, size=9, color=DIM, font="Calibri")

footer(s, 15)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — TEAM & RBAC SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "TEAM & RBAC SYSTEM", "Multi-team isolation and role-based access control")

# Role matrix
card(s, 0.4, 1.05, 12.5, 2.3, title="ROLE CAPABILITIES MATRIX")
roles_h = ["CAPABILITY", "SUPERADMIN", "ADMIN", "ANALYST", "VIEWER"]
col_xs  = [0.6, 3.5, 6.0, 8.5, 10.8]
role_cols = [WHITE, YELLOW, CYAN, GREEN, DIM]
for i, (h, col) in enumerate(zip(roles_h, role_cols)):
    txt(s, h, col_xs[i], 1.42, 2.3, 0.28, size=8.5, bold=True, color=col, font="Calibri")

matrix = [
    ("Create/delete teams",     "✓", "✗", "✗", "✗"),
    ("Add any user to any team","✓", "✗", "✗", "✗"),
    ("Add member to own team",  "✓", "✓", "✗", "✗"),
    ("Invite new users",        "✓", "✓", "✗", "✗"),
    ("See all accounts",        "✓", "own team", "own team", "own team"),
    ("Add cloud account",       "✓", "✓", "✓", "✗"),
    ("Delete cloud account",    "✓", "✓", "✗", "✗"),
    ("Scan account",            "✓", "✓", "✓", "✗"),
    ("Create custom policies",  "✓", "✓", "✓", "✗"),
    ("View audit log",          "✓", "✓", "own only", "✗"),
]
for ri, (cap, *vals) in enumerate(matrix):
    ry = 1.72 + ri * 0.16
    txt(s, cap, col_xs[0], ry, 2.7, 0.18, size=7.5, color=DIM, font="Calibri")
    for ci, val in enumerate(vals):
        col = GREEN if val == "✓" else (RED if val == "✗" else ORANGE)
        txt(s, val, col_xs[ci+1], ry, 2.2, 0.18, size=7.5, bold=(val in "✓✗"), color=col, font="Calibri")

# Team structure
card(s, 0.4, 3.55, 5.9, 3.35, title="TEAM STRUCTURE")
team_pts = [
    "Teams are isolated silos — accounts are scoped to one team",
    "Users can belong to multiple teams (many-to-many)",
    "team_members table: team_id × user_id with UNIQUE constraint",
    "cloud_accounts.team_id FK — NULL = superadmin-only visible",
    "Deleting a team: accounts get team_id = NULL (not deleted)",
    "Admin auto-added to new user's team on invite-accept",
    "Superadmin sees and manages all teams from one panel",
    "Team member additions take effect on next API call",
]
bullet_block(s, team_pts, 0.6, 3.95, w=5.5, size=9.5, dot_color=CYAN)

# DB schema
card(s, 6.55, 3.55, 6.35, 3.35, title="DATABASE SCHEMA — TEAMS")
txt(s,
    'teams\n'
    '  id UUID PK\n'
    '  name TEXT UNIQUE\n'
    '  description TEXT\n'
    '  created_by UUID → users\n'
    '  created_at TIMESTAMPTZ\n\n'
    'team_members\n'
    '  id UUID PK\n'
    '  team_id UUID → teams (CASCADE)\n'
    '  user_id UUID → users (CASCADE)\n'
    '  added_by UUID → users\n'
    '  UNIQUE (team_id, user_id)',
    6.75, 3.92, 5.9, 2.8, size=8.5, color=GREEN, font="Consolas")

footer(s, 16)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — DATABASE SCHEMA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "DATABASE SCHEMA", "PostgreSQL 15 — 11 tables")

tables = [
    ("users", YELLOW, [
        "id UUID PK",
        "email, username",
        "password_hash (bcrypt)",
        "role (superadmin/admin/analyst/viewer)",
        "is_active, mfa_enabled",
        "created_by UUID → users",
    ]),
    ("cloud_accounts", CYAN, [
        "id UUID PK",
        "name, cloud (aws/azure)",
        "encrypted_creds (Fernet)",
        "team_id UUID → teams",
        "last_scanned_at",
        "scan_interval_hours",
    ]),
    ("scan_results", GREEN, [
        "id UUID PK",
        "account_id UUID → cloud_accounts",
        "findings JSONB (full list)",
        "risk_score, resource_count",
        "critical/high/medium/low counts",
        "scanned_at TIMESTAMPTZ",
    ]),
    ("teams / team_members", ORANGE, [
        "teams: id, name, description",
        "team_members: team_id × user_id",
        "Both with CASCADE deletes",
        "Unique constraint on membership",
        "Indexed on team_id, user_id",
    ]),
    ("alert_settings", RED, [
        "id UUID PK",
        "user_id, account_id",
        "threshold_score (int)",
        "alert_on_critical (bool)",
        "recipient_email",
        "enabled (bool)",
    ]),
    ("audit_log", DIM, [
        "id UUID PK",
        "user_id, user_email",
        "action (event key)",
        "resource_id, resource_name",
        "ip_address",
        "detail JSONB",
    ]),
]

for i, (name, col, fields) in enumerate(tables):
    col_n = i % 3
    row_n = i // 3
    cx = 0.4 + col_n * 4.3
    cy = 1.1 + row_n * 3.15
    card(s, cx, cy, 4.05, 3.0, title=name, title_color=col)
    bullet_block(s, fields, cx + 0.2, cy + 0.45, w=3.65, size=9.5, dot_color=col, line_gap=0.34)

# Extra tables note
txt(s, "Additional tables: invite_tokens  ·  password_reset_tokens  ·  finding_statuses  ·  "
       "scheduled_jobs  ·  alert_history  ·  system_alert_settings",
    0.4, 7.05, 12.8, 0.25, size=8, color=DIM, font="Calibri")

footer(s, 17)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "DEPLOYMENT", "Docker Compose, Nginx, and production configuration")

# Services
card(s, 0.4, 1.05, 6.0, 3.55, title="DOCKER COMPOSE SERVICES")
services = [
    ("db",       "postgres:15", "PostgreSQL with health check, volume mount", CYAN),
    ("backend",  "python:3.11", "FastAPI + uvicorn on port 8000", YELLOW),
    ("frontend", "node:20",     "Vite dev server (dev) / nginx (prod)", GREEN),
    ("nginx",    "nginx:alpine","Reverse proxy, TLS termination, port 80/443", ORANGE),
]
for i, (svc, img, desc, col) in enumerate(services):
    sy = 1.45 + i * 0.78
    txt(s, svc,  0.6, sy,       1.3, 0.28, size=10, bold=True, color=col, font="Consolas")
    txt(s, img,  2.0, sy,       2.2, 0.28, size=9,             color=DIM, font="Calibri")
    txt(s, desc, 0.6, sy+0.3,   5.6, 0.35, size=9,             color=DIM, font="Calibri")

# Environment
card(s, 6.65, 1.05, 6.25, 3.55, title=".ENV CONFIGURATION")
env = [
    "DATABASE_URL — asyncpg connection string",
    "SECRET_KEY — JWT signing secret",
    "FERNET_KEY — credential encryption key",
    "SMTP_HOST / SMTP_PORT / SMTP_USER / SMTP_PASS",
    "VITE_API_URL — frontend API base URL",
    "CORS_ORIGINS — allowed frontend origins",
    "FIRST_ADMIN_EMAIL / PASSWORD (setup mode)",
]
bullet_block(s, env, 6.85, 1.45, w=5.85, size=9.5, dot_color=CYAN)

# Dev vs Prod
card(s, 0.4, 4.8, 5.9, 2.1, title="DEVELOPMENT vs PRODUCTION")
txt(s, "DEV",  0.6, 5.15, 0.5, 0.25, size=9, bold=True, color=YELLOW, font="Calibri")
txt(s, "docker-compose.yml — hot-reload Vite + FastAPI, debug logs, no TLS",
    0.6, 5.15, 5.5, 0.28, size=9, color=DIM, font="Calibri")
txt(s, "PROD", 0.6, 5.6, 0.55, 0.25, size=9, bold=True, color=GREEN, font="Calibri")
txt(s, "docker-compose.prod.yml — static frontend build, Nginx TLS, uvicorn workers",
    0.6, 5.6, 5.5, 0.28, size=9, color=DIM, font="Calibri")

# Nginx config summary
card(s, 6.65, 4.8, 6.25, 2.1, title="NGINX ROUTING")
nginx = [
    "/api/* → http://backend:8000 (proxy_pass)",
    "/ → static build or frontend:5173 (dev)",
    "CORS headers set at proxy layer",
    "WebSocket upgrade headers (for future WS)",
    "Gzip compression enabled for JS/CSS",
]
bullet_block(s, nginx, 6.85, 5.15, w=5.85, size=9.5, dot_color=ORANGE)

footer(s, 18)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — SECURITY SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "SECURITY DESIGN SUMMARY", "Defense-in-depth across all layers")

pillars = [
    ("Authentication", YELLOW, [
        "bcrypt (cost 12) password hashing",
        "JWT HS256 tokens, 24h TTL",
        "TOTP MFA (RFC 6238) with QR provisioning",
        "Signed single-use password-reset tokens (1h TTL)",
        "Account active/inactive flag with immediate lockout",
    ]),
    ("Authorization", CYAN, [
        "require_role() FastAPI dependency on every endpoint",
        "4-level role hierarchy: superadmin > admin > analyst > viewer",
        "Team isolation: data filtered by team membership at DB query level",
        "Admin cannot escalate to admin/superadmin roles",
        "Admin RBAC scoped: can only manage own-team lower-role users",
    ]),
    ("Data Protection", GREEN, [
        "Fernet (AES-128-CBC + HMAC-SHA256) for cloud credentials",
        "Encryption key in server env — never in database",
        "Credentials never returned in any API response",
        "PostgreSQL TLS connection (production)",
        "HTTPS-only in production (Nginx TLS termination)",
    ]),
    ("Audit & Monitoring", ORANGE, [
        "Every significant action logged to audit_log table",
        "IP address captured from X-Forwarded-For / direct",
        "Detail JSONB captures context (old role, new role, etc.)",
        "Immutable log — no edit/delete endpoints for audit entries",
        "Superadmin has full platform-wide visibility",
    ]),
]

for i, (pillar, col, items) in enumerate(pillars):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.05, 3.05, 5.85, title=pillar, title_color=col)
    bullet_block(s, items, cx + 0.18, 1.45, w=2.7, size=9.5, dot_color=col)

footer(s, 19)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — SUMMARY & ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "SUMMARY & FUTURE ROADMAP", "What's built and what comes next")

# Built summary
card(s, 0.4, 1.05, 6.0, 5.85, title="WHAT'S BUILT")
done = [
    "Multi-cloud CSPM for AWS (90+ rules) and Azure (137 rules)",
    "Full compliance mapping: CIS, NIST, PCI, HIPAA, SOC2, ISO, GDPR",
    "Real-time and scheduled scanning with APScheduler",
    "Risk scoring: 0–100 penalty model with severity weighting",
    "Team-based RBAC: superadmin / admin / analyst / viewer",
    "Team isolation: accounts and findings scoped per team",
    "TOTP MFA with QR code provisioning (pyotp)",
    "Fernet-encrypted credential storage",
    "Full audit log with IP tracking and JSON detail",
    "Custom policy engine (10 operators, any resource_type)",
    "Email alerting: HTML reports, score threshold, critical-finding triggers",
    "Scheduled background scans (per-account interval)",
    "Finding status tracking: open / acknowledged / resolved",
    "First-time setup wizard (superadmin bootstrap)",
    "Docker Compose dev + prod, Nginx reverse proxy",
    "Dark/light theme, fully custom React UI (no external UI lib)",
]
bullet_block(s, done, 0.6, 1.45, w=5.5, size=9.5, dot_color=GREEN)

# Roadmap
card(s, 6.65, 1.05, 6.25, 5.85, title="FUTURE ROADMAP")
roadmap_sections = [
    ("Short Term", YELLOW, [
        "GCP connector (Cloud SQL, GCS, IAM, GKE)",
        "Expand AWS rules to 220+ (Redshift, WAF, Cognito, ECS…)",
        "Unit + integration test coverage",
        "Nginx prod config — serve static build correctly",
    ]),
    ("Medium Term", CYAN, [
        "SAML / SSO integration (Okta, Azure AD)",
        "Webhook alerts (Slack, Teams, PagerDuty)",
        "CSV / PDF compliance report export",
        "Role-specific dashboard views",
        "Resource-level history and diff view",
    ]),
    ("Long Term", ORANGE, [
        "Auto-remediation with approval workflows",
        "Terraform / IaC drift detection",
        "CI/CD pipeline security gate integration",
        "Multi-region deployment with data residency controls",
        "AI-powered remediation suggestions (LLM-assisted)",
        "SOC 2 Type II evidence collection automation",
    ]),
]
ry = 1.45
for section, col, items in roadmap_sections:
    txt(s, section, 6.85, ry, 5.5, 0.28, size=10, bold=True, color=col, font="Calibri")
    bullet_block(s, items, 6.85, ry + 0.28, w=5.85, size=9.5, dot_color=col)
    ry += 0.28 + len(items) * 0.28 + 0.22

footer(s, 20)


# ─── Save ────────────────────────────────────────────────────────────────────
out = "/Users/anishdoulagar/Documents/CSPM/Vanguard_CSPM_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
